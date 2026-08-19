"""Build Lesson 01 PPTX from slides.md (32 pages, 16:9, all native shapes)."""
import os
import re
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# --------- paths ---------
SLIDES_MD = "/Users/nathan/Projects/tc/01-ai-agent-engineering/slides/lesson-01-slides.md"
OUT = "/Users/nathan/Projects/tc/01-ai-agent-engineering/ppt/AI Agent系统工程实战 - 第1课：从聊天机器人到智能助手.pptx"
ASSETS = "/Users/nathan/Projects/tc/01-ai-agent-engineering/ppt/assets"
COVER_IMG = os.path.join(ASSETS, "cover.png")
CLOSING_IMG = os.path.join(ASSETS, "closing.png")

# --------- design tokens ---------
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

NAVY = RGBColor(0x0E, 0x2A, 0x47)         # 主深蓝
BLUE = RGBColor(0x1E, 0x4F, 0x8B)         # 中蓝
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFA)   # 淡蓝底
ACCENT = RGBColor(0xF5, 0x8E, 0x2C)       # 强调橙
ACCENT_DARK = RGBColor(0xC9, 0x6B, 0x14)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x6F, 0x7A)
GRAY_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG = RGBColor(0xE6, 0xE6, 0xE6)
CODE_KEY = RGBColor(0xFF, 0xB7, 0x6B)
CODE_STR = RGBColor(0xA8, 0xE0, 0xA8)
CODE_COMMENT = RGBColor(0x88, 0x88, 0xAA)

CN_FONT = "微软雅黑"
CODE_FONT = "Consolas"


# --------- parser ---------
def parse_slides(path):
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()
    # split by horizontal rule that delimits pages: lines of just '---' on its own
    parts = re.split(r"\n---\n", text)
    pages = []
    for part in parts:
        part = part.strip()
        if not part:
            continue
        # find a header like ## [TYPE] Slide ...
        m = re.search(r"^##\s*\[([A-Z]+)\]\s*(.+)$", part, re.MULTILINE)
        if not m:
            continue
        ptype = m.group(1)
        slide_id = m.group(2).strip()
        # body = everything after that header line
        body = part[m.end():].strip()
        # extract notes (lines starting with > 对应讲课稿)
        notes = []
        body_lines = []
        for ln in body.split("\n"):
            if ln.strip().startswith(">"):
                notes.append(ln.strip().lstrip(">").strip())
            else:
                body_lines.append(ln)
        body_clean = "\n".join(body_lines).strip()
        pages.append({
            "type": ptype,
            "id": slide_id,
            "body": body_clean,
            "notes": "\n".join(notes),
        })
    return pages


# --------- helpers ---------
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    if not shadow:
        # disable default shadow
        sppr = shape.shadow._element  # type: ignore
    return shape


def add_round_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=GRAY_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=CN_FONT, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        # ensure east asian font
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font)
    return tb


def add_rich_paragraphs(slide, x, y, w, h, paragraphs, *, anchor=MSO_ANCHOR.TOP):
    """paragraphs: list of {runs:[(text,opts)], align, space_after}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if "line_spacing" in para:
            p.line_spacing = para["line_spacing"]
        if "space_after" in para:
            p.space_after = Pt(para["space_after"])
        for text, opts in para["runs"]:
            run = p.add_run()
            run.text = text
            run.font.name = opts.get("font", CN_FONT)
            run.font.size = Pt(opts.get("size", 16))
            run.font.bold = opts.get("bold", False)
            run.font.italic = opts.get("italic", False)
            run.font.color.rgb = opts.get("color", GRAY_DARK)
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn("a:ea"))
            if ea is None:
                ea = etree.SubElement(rPr, qn("a:ea"))
            ea.set("typeface", opts.get("font", CN_FONT))
    return tb


def add_arrow_down(slide, x_center, y_top, length, color=BLUE, width_pt=2.5):
    """Vertical down arrow as a connector."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x_center, y_top, x_center, y_top + length)
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    # add arrow head
    line = conn.line._get_or_add_ln()
    tail = etree.SubElement(line, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return conn


def add_arrow_right(slide, x, y, length, color=BLUE, width_pt=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + length, y)
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    line = conn.line._get_or_add_ln()
    tail = etree.SubElement(line, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return conn


def add_page_header(slide, title, idx, total, subtitle=None):
    # top accent bar
    bar = add_rect(slide, Emu(0), Emu(0), SLIDE_W, Emu(100000), NAVY)
    add_rect(slide, Emu(0), Emu(100000), Emu(2200000), Emu(50000), ACCENT)
    # title
    add_text(slide, Emu(600000), Emu(280000), Emu(9000000), Emu(700000),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Emu(600000), Emu(900000), Emu(9000000), Emu(450000),
                 subtitle, size=16, color=GRAY)
    # page number
    add_text(slide, Emu(10800000), Emu(6400000), Emu(1300000), Emu(350000),
             f"{idx} / {total}", size=11, color=GRAY, align=PP_ALIGN.RIGHT)
    # course mark
    add_text(slide, Emu(400000), Emu(6400000), Emu(6000000), Emu(350000),
             "AI Agent系统工程实战 · 第1课 · 记忆系统", size=10, color=GRAY)


def add_notes(slide, notes):
    if not notes:
        return
    nslide = slide.notes_slide
    tf = nslide.notes_text_frame
    tf.text = notes


# --------- code highlight (very lightweight) ---------
PY_KEYWORDS = {"def", "class", "return", "if", "else", "elif", "import", "from",
               "for", "in", "while", "self", "True", "False", "None", "pass", "lambda",
               "and", "or", "not", "as", "with", "try", "except", "finally", "raise"}


def render_code_lines(tf, code, *, size=14):
    tf.word_wrap = False
    tf.margin_left = Emu(180000); tf.margin_right = Emu(180000)
    tf.margin_top = Emu(140000); tf.margin_bottom = Emu(140000)
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        # split into tokens by simple regex
        # handle inline comment
        comment_idx = line.find("#")
        in_str = False
        quote = None
        # naive split: token = word|space|punct
        seg = line
        comment = None
        # detect '#' not inside string
        cur_in_str = False
        cur_q = None
        cidx = -1
        for k, ch in enumerate(line):
            if cur_in_str:
                if ch == cur_q:
                    cur_in_str = False
                continue
            if ch in ('"', "'"):
                cur_in_str = True
                cur_q = ch
                continue
            if ch == "#":
                cidx = k
                break
        if cidx >= 0:
            seg = line[:cidx]
            comment = line[cidx:]
        # tokenize seg
        tokens = re.findall(r'(\s+|"[^"]*"|\'[^\']*\'|[A-Za-z_][A-Za-z_0-9]*|\d+|.)', seg)
        for tok in tokens:
            run = p.add_run()
            run.text = tok
            run.font.name = CODE_FONT
            run.font.size = Pt(size)
            if tok.strip() == "":
                run.font.color.rgb = CODE_FG
            elif (tok.startswith('"') and tok.endswith('"')) or (tok.startswith("'") and tok.endswith("'")):
                run.font.color.rgb = CODE_STR
            elif tok in PY_KEYWORDS:
                run.font.color.rgb = CODE_KEY
                run.font.bold = True
            elif re.match(r"^\d+$", tok):
                run.font.color.rgb = CODE_STR
            else:
                run.font.color.rgb = CODE_FG
        if comment is not None:
            run = p.add_run()
            run.text = comment
            run.font.name = CODE_FONT
            run.font.size = Pt(size)
            run.font.color.rgb = CODE_COMMENT
            run.font.italic = True


# ===================== RENDERERS =====================

def render_cover(slide, page, idx, total):
    set_slide_bg(slide, NAVY)
    # background image (right half) if present
    if os.path.exists(COVER_IMG):
        try:
            slide.shapes.add_picture(COVER_IMG, Emu(0), Emu(0), width=SLIDE_W, height=SLIDE_H)
            # dark overlay for text legibility
            overlay = add_rect(slide, Emu(0), Emu(0), Emu(7400000), SLIDE_H, NAVY)
            overlay.fill.transparency = 0  # solid
            # Actually use semi-transparent via XML
            sp = overlay.fill.fore_color._xFill
            # add alpha via transparency element
            solidFill = sp
            srgb = solidFill.find(qn("a:srgbClr"))
            if srgb is not None:
                alpha = etree.SubElement(srgb, qn("a:alpha"))
                alpha.set("val", "75000")
        except Exception as e:
            print("cover image embed failed:", e)
    # accent bar
    add_rect(slide, Emu(600000), Emu(2400000), Emu(120000), Emu(900000), ACCENT)
    # course series
    add_text(slide, Emu(800000), Emu(2350000), Emu(6000000), Emu(500000),
             "AI Agent 系统工程实战", size=22, bold=True, color=ACCENT)
    # main title
    add_text(slide, Emu(800000), Emu(2900000), Emu(6800000), Emu(900000),
             "第1课：从聊天机器人到智能助手", size=36, bold=True, color=WHITE)
    add_text(slide, Emu(800000), Emu(3700000), Emu(6800000), Emu(700000),
             "—— AI Agent 的记忆觉醒", size=26, bold=True, color=WHITE)
    # meta line
    add_text(slide, Emu(800000), Emu(4900000), Emu(6800000), Emu(500000),
             "对应原文：第1-2章   |   难度：中级   |   时长：3小时",
             size=14, color=GRAY_LIGHT)
    # bottom strip
    add_rect(slide, Emu(0), Emu(6500000), SLIDE_W, Emu(358000), NAVY)
    add_text(slide, Emu(600000), Emu(6560000), Emu(8000000), Emu(280000),
             "讲义：lesson-01-memory-system.md   ·   讲稿：lesson-01-script.md",
             size=10, color=GRAY_LIGHT)
    add_text(slide, Emu(10000000), Emu(6560000), Emu(1700000), Emu(280000),
             f"{idx} / {total}", size=10, color=GRAY_LIGHT, align=PP_ALIGN.RIGHT)
    add_notes(slide, page["notes"])


def render_toc(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    add_page_header(slide, "今天的旅程", idx, total, subtitle="Course Roadmap · 5 modules + 2 exercises")
    # parse table from body
    lines = [ln for ln in page["body"].split("\n") if ln.strip().startswith("|")]
    rows = []
    for ln in lines:
        cells = [c.strip() for c in ln.strip().strip("|").split("|")]
        rows.append(cells)
    # rows[0]=header, rows[1]=separator (---), rows[2:]=data
    data = [r for r in rows[2:] if r and not all(set(c) <= set("-") for c in r)]
    # render as a styled table using shapes (not pptx table for stronger styling) - use pptx table
    from pptx.util import Cm
    rows_n = len(data) + 1
    cols_n = 3
    tbl_x = Emu(700000); tbl_y = Emu(1500000)
    tbl_w = Emu(10800000); tbl_h = Emu(4400000)
    table_shape = slide.shapes.add_table(rows_n, cols_n, tbl_x, tbl_y, tbl_w, tbl_h)
    table = table_shape.table
    table.columns[0].width = Emu(1300000)
    table.columns[1].width = Emu(7500000)
    table.columns[2].width = Emu(2000000)
    headers = ["模块", "主题", "时间"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.name = CN_FONT; run.font.size = Pt(18); run.font.bold = True
        run.font.color.rgb = WHITE
    for i, row in enumerate(data, start=1):
        for j, txt in enumerate(row[:3]):
            cell = table.cell(i, j)
            cell.text = ""
            is_break = txt.startswith("--") or "练习" in txt or "休息" in txt or "总结" in txt
            cell.fill.solid()
            if is_break:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
            else:
                cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT_BLUE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT
            run = p.add_run(); run.text = txt.replace("--", "—")
            run.font.name = CN_FONT
            run.font.size = Pt(15)
            run.font.bold = (j == 0) or is_break
            run.font.color.rgb = ACCENT_DARK if is_break else NAVY if j == 0 else GRAY_DARK
    add_notes(slide, page["notes"])


def render_question(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    add_page_header(slide, "互动思考", idx, total, subtitle="Question · Think for 30 seconds")
    body = page["body"]
    # find first H1 = main question
    m = re.search(r"^#\s+(.+)$", body, re.MULTILINE)
    main_q = m.group(1).strip() if m else ""
    # rest = body without that line
    rest = body
    if m:
        rest = body[:m.start()] + body[m.end():]
    rest = rest.strip()

    # large question card
    card = add_round_rect(slide, Emu(800000), Emu(1700000), Emu(10600000), Emu(1500000),
                          LIGHT_BLUE, line_color=BLUE)
    card.line.width = Pt(2)
    # question mark icon
    qmark = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(950000), Emu(1900000), Emu(900000), Emu(900000))
    qmark.fill.solid(); qmark.fill.fore_color.rgb = ACCENT
    qmark.line.fill.background()
    add_text(slide, Emu(950000), Emu(1900000), Emu(900000), Emu(900000),
             "?", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Emu(2050000), Emu(2050000), Emu(9300000), Emu(1100000),
             main_q, size=28, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    # rest content - render lines, options as choices
    y = Emu(3500000)
    lines = [ln for ln in rest.split("\n") if ln.strip()]
    option_pat = re.compile(r"^-\s*([A-D])\.\s*(.+)$")
    options = []
    other = []
    for ln in lines:
        m2 = option_pat.match(ln.strip())
        if m2:
            options.append((m2.group(1), m2.group(2)))
        else:
            other.append(ln.strip().lstrip("- ").strip())
    if options:
        opt_w = Emu(3400000); opt_h = Emu(1500000); gap = Emu(200000)
        total_w = opt_w * len(options) + gap * (len(options) - 1)
        x0 = (SLIDE_W - total_w) // 2
        for i, (label, txt) in enumerate(options):
            x = x0 + i * (opt_w + gap)
            box = add_round_rect(slide, x, y, opt_w, opt_h, WHITE, line_color=BLUE)
            box.line.width = Pt(1.5)
            add_text(slide, x, y + Emu(100000), opt_w, Emu(500000),
                     label, size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
            add_text(slide, x + Emu(200000), y + Emu(700000), opt_w - Emu(400000), Emu(700000),
                     txt.replace("**", ""), size=15, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    if other:
        y2 = y + Emu(1700000) if options else y
        add_text(slide, Emu(800000), y2, Emu(10600000), Emu(800000),
                 "\n".join(o.replace("**", "") for o in other),
                 size=18, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=1.4)
    add_notes(slide, page["notes"])


def render_concept(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    body = page["body"]
    # find module subtitle ### ...
    sub_m = re.search(r"^###\s+(.+)$", body, re.MULTILINE)
    subtitle = sub_m.group(1).strip() if sub_m else None
    # main title from first # ...
    title_m = re.search(r"^#\s+(.+)$", body, re.MULTILINE)
    title = title_m.group(1).strip() if title_m else page["id"]

    add_page_header(slide, title, idx, total, subtitle=subtitle)

    # remove header/title/subtitle lines
    rest = body
    if sub_m: rest = rest.replace(sub_m.group(0), "")
    if title_m: rest = rest.replace(title_m.group(0), "")
    rest = rest.strip()

    y = Emu(1700000)
    # detect bullet list
    bullet_lines = []
    paragraphs = []
    for ln in rest.split("\n"):
        s = ln.strip()
        if not s: continue
        if s.startswith("- "):
            bullet_lines.append(s[2:].strip())
        else:
            paragraphs.append(s)

    # render paragraphs as highlight blocks
    for p in paragraphs:
        is_quote = p.startswith(">")
        ptxt = p.lstrip(">").strip()
        # split bold **...**
        parts = re.split(r"(\*\*[^*]+\*\*)", ptxt)
        runs = []
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                runs.append((part[2:-2], {"size": 22, "bold": True, "color": ACCENT_DARK}))
            else:
                runs.append((part, {"size": 22, "bold": False, "color": NAVY if not is_quote else GRAY}))
        # quote box
        if is_quote:
            add_round_rect(slide, Emu(700000), y, Emu(10800000), Emu(700000), LIGHT_BLUE)
            add_rect(slide, Emu(700000), y, Emu(80000), Emu(700000), ACCENT)
            add_rich_paragraphs(slide, Emu(900000), y + Emu(100000), Emu(10500000), Emu(500000),
                                [{"runs": runs, "align": PP_ALIGN.LEFT}],
                                anchor=MSO_ANCHOR.MIDDLE)
            y += Emu(900000)
        else:
            add_rich_paragraphs(slide, Emu(700000), y, Emu(10800000), Emu(700000),
                                [{"runs": runs, "align": PP_ALIGN.LEFT, "line_spacing": 1.4}])
            y += Emu(800000)

    # bullet list rendered as cards
    if bullet_lines:
        n = len(bullet_lines)
        if n <= 4:
            card_h = Emu(700000); gap = Emu(140000)
            for i, b in enumerate(bullet_lines):
                yy = y + i * (card_h + gap)
                # split "label = value" pattern
                if "=" in b:
                    left, right = b.split("=", 1)
                    add_round_rect(slide, Emu(700000), yy, Emu(10800000), card_h, WHITE, line_color=BLUE)
                    add_rect(slide, Emu(700000), yy, Emu(120000), card_h, ACCENT)
                    runs = [
                        (left.strip().replace("**", ""), {"size": 18, "bold": True, "color": NAVY}),
                        ("  =  ", {"size": 18, "color": GRAY}),
                        (right.strip().replace("**", ""), {"size": 18, "color": GRAY_DARK}),
                    ]
                    add_rich_paragraphs(slide, Emu(900000), yy, Emu(10500000), card_h,
                                        [{"runs": runs, "align": PP_ALIGN.LEFT}],
                                        anchor=MSO_ANCHOR.MIDDLE)
                else:
                    add_round_rect(slide, Emu(700000), yy, Emu(10800000), card_h, LIGHT_BLUE)
                    add_text(slide, Emu(900000), yy, Emu(10500000), card_h,
                             "• " + b.replace("**", ""), size=18, color=NAVY,
                             anchor=MSO_ANCHOR.MIDDLE)
        else:
            for i, b in enumerate(bullet_lines):
                add_text(slide, Emu(800000), y + i * Emu(500000), Emu(10500000), Emu(500000),
                         "• " + b.replace("**", ""), size=18, color=GRAY_DARK)
    add_notes(slide, page["notes"])


# helper for diagram pages
def render_data_flow(slide, x, y, items, *, box_w=Emu(5000000), box_h=Emu(550000),
                     gap=Emu(180000), arrow_pad=Emu(40000)):
    """Vertical flow: list of (text, [optional sub_text]). Returns final y."""
    cx = x + box_w // 2
    cur_y = y
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, sub = item
        else:
            text, sub = item, None
        box = add_round_rect(slide, x, cur_y, box_w, box_h, LIGHT_BLUE, line_color=BLUE)
        box.line.width = Pt(1.5)
        add_text(slide, x, cur_y, box_w, box_h, text, size=16, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if sub:
            add_text(slide, x + box_w + Emu(200000), cur_y, Emu(5000000), box_h, sub,
                     size=13, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
        cur_y += box_h
        if i != len(items) - 1:
            add_arrow_down(slide, cx, cur_y + arrow_pad, gap - 2 * arrow_pad)
            cur_y += gap
    return cur_y


def render_diagram_clawbot_v1(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    add_page_header(slide, "ClawBot v1 数据流", idx, total, subtitle="Module 1.1 · 四步工作流")
    items = [
        ("用户消息", "User Message"),
        ("加载 memory.json", "全部历史对话"),
        ("构建 messages 列表", "system + history[-20:] + 当前消息"),
        ("调用 GPT-4 API", "openai.ChatCompletion.create()"),
        ("保存到 memory.json", "只保留最近 20 轮"),
        ("返回回复", "Response"),
    ]
    render_data_flow(slide, Emu(2400000), Emu(1500000), items,
                     box_w=Emu(4500000), box_h=Emu(530000), gap=Emu(170000))
    # callout: bottleneck
    cx = Emu(7400000); cy = Emu(1500000)
    add_round_rect(slide, cx, cy, Emu(4400000), Emu(4500000), RGBColor(0xFF, 0xF3, 0xE0),
                   line_color=ACCENT)
    add_text(slide, cx, cy + Emu(150000), Emu(4400000), Emu(500000),
             "⚠ 关键瓶颈", size=20, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, cx + Emu(300000), cy + Emu(800000), Emu(3800000), Emu(3500000),
             "history[-20:]\n一刀切丢弃旧信息\n\n• 第 21 轮起永久遗忘\n• 不区分重要性\n• 跨会话无检索\n• 无语义关联",
             size=15, color=GRAY_DARK, line_spacing=1.5)
    add_notes(slide, page["notes"])


def render_diagram_embedding(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    add_page_header(slide, "Embedding 的直觉理解", idx, total,
                    subtitle="语义相近 → 空间距离更近")
    # plot area
    px = Emu(800000); py = Emu(1600000)
    pw = Emu(6500000); ph = Emu(4500000)
    add_round_rect(slide, px, py, pw, ph, RGBColor(0xFA, 0xFA, 0xFA), line_color=GRAY_LIGHT)
    # axes
    ax_color = GRAY
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px + Emu(400000), py + ph - Emu(400000),
                               px + pw - Emu(200000), py + ph - Emu(400000)).line.color.rgb = ax_color
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px + Emu(400000), py + ph - Emu(400000),
                               px + Emu(400000), py + Emu(200000)).line.color.rgb = ax_color
    add_text(slide, px + pw - Emu(800000), py + ph - Emu(380000), Emu(700000), Emu(300000),
             "维度 1", size=10, color=GRAY)
    add_text(slide, px + Emu(420000), py + Emu(80000), Emu(700000), Emu(300000),
             "维度 2", size=10, color=GRAY)
    # points: (rel_x, rel_y, label, color)
    points = [
        (0.20, 0.85, "男人", BLUE),
        (0.80, 0.75, "国王", BLUE),
        (0.30, 0.30, "女人", ACCENT),
        (0.85, 0.20, "女王", ACCENT),
    ]
    for rx, ry, lbl, col in points:
        cx = px + Emu(int(pw * rx))
        cy = py + Emu(int(ph * ry))
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Emu(120000), cy - Emu(120000),
                                     Emu(240000), Emu(240000))
        dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
        add_text(slide, cx + Emu(150000), cy - Emu(160000), Emu(1500000), Emu(320000),
                 lbl, size=14, bold=True, color=col)
    # vector arrows demonstrating arithmetic
    def line(x1, y1, x2, y2, color, dash=False):
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
            px + Emu(int(pw*x1)), py + Emu(int(ph*y1)),
            px + Emu(int(pw*x2)), py + Emu(int(ph*y2)))
        c.line.color.rgb = color
        c.line.width = Pt(1.5)
        if dash:
            ln = c.line._get_or_add_ln()
            prstDash = etree.SubElement(ln, qn("a:prstDash"))
            prstDash.set("val", "dash")
        return c
    line(0.20, 0.85, 0.80, 0.75, BLUE, dash=True)
    line(0.30, 0.30, 0.85, 0.20, ACCENT, dash=True)

    # right side: equation card
    rx = Emu(7600000); ry = Emu(1700000)
    add_round_rect(slide, rx, ry, Emu(4200000), Emu(1500000), LIGHT_BLUE, line_color=BLUE)
    add_text(slide, rx, ry + Emu(150000), Emu(4200000), Emu(450000),
             "向量算术", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, rx + Emu(200000), ry + Emu(700000), Emu(3800000), Emu(700000),
             "国王 − 男人 + 女人 ≈ 女王", size=20, bold=True, color=ACCENT_DARK,
             align=PP_ALIGN.CENTER, font=CODE_FONT)
    # synonym card
    add_round_rect(slide, rx, ry + Emu(1700000), Emu(4200000), Emu(2500000),
                   RGBColor(0xFF, 0xF7, 0xEC), line_color=ACCENT)
    add_text(slide, rx, ry + Emu(1850000), Emu(4200000), Emu(450000),
             "同义词自动聚拢", size=18, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER)
    syn = "API密钥  ≈  API key  ≈  访问令牌\n\n交易所密钥  ≈  binance API key"
    add_text(slide, rx + Emu(200000), ry + Emu(2400000), Emu(3800000), Emu(1700000),
             syn, size=14, color=GRAY_DARK, align=PP_ALIGN.CENTER, line_spacing=1.6)
    add_notes(slide, page["notes"])


def render_diagram_cosine(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    add_page_header(slide, "余弦相似度", idx, total,
                    subtitle="衡量两个向量方向的接近程度")
    # formula card
    fx = Emu(800000); fy = Emu(1600000)
    add_round_rect(slide, fx, fy, Emu(10600000), Emu(900000), NAVY)
    add_text(slide, fx, fy + Emu(150000), Emu(10600000), Emu(300000),
             "公式", size=14, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, fx, fy + Emu(450000), Emu(10600000), Emu(450000),
             "cosine(A, B) = (A · B) / (|A| × |B|)    →    结果范围 0 ~ 1",
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=CODE_FONT)
    # examples - two cards side by side
    sim_y = Emu(2750000)
    # high similarity
    add_round_rect(slide, Emu(800000), sim_y, Emu(5100000), Emu(1900000),
                   RGBColor(0xE6, 0xF7, 0xE6), line_color=GREEN)
    add_text(slide, Emu(800000), sim_y + Emu(180000), Emu(5100000), Emu(400000),
             "✓ 高度相似", size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, Emu(900000), sim_y + Emu(700000), Emu(4900000), Emu(400000),
             'cosine("API密钥", "API key")', size=14, color=GRAY_DARK,
             align=PP_ALIGN.CENTER, font=CODE_FONT)
    add_text(slide, Emu(900000), sim_y + Emu(1200000), Emu(4900000), Emu(600000),
             "≈ 0.92", size=36, bold=True, color=GREEN, align=PP_ALIGN.CENTER, font=CODE_FONT)
    # low similarity
    add_round_rect(slide, Emu(6300000), sim_y, Emu(5100000), Emu(1900000),
                   RGBColor(0xFD, 0xEA, 0xEA), line_color=RGBColor(0xC0, 0x39, 0x2B))
    add_text(slide, Emu(6300000), sim_y + Emu(180000), Emu(5100000), Emu(400000),
             "✗ 几乎无关", size=18, bold=True,
             color=RGBColor(0xC0, 0x39, 0x2B), align=PP_ALIGN.CENTER)
    add_text(slide, Emu(6400000), sim_y + Emu(700000), Emu(4900000), Emu(400000),
             'cosine("API密钥", "天气预报")', size=14, color=GRAY_DARK,
             align=PP_ALIGN.CENTER, font=CODE_FONT)
    add_text(slide, Emu(6400000), sim_y + Emu(1200000), Emu(4900000), Emu(600000),
             "≈ 0.05", size=36, bold=True,
             color=RGBColor(0xC0, 0x39, 0x2B), align=PP_ALIGN.CENTER, font=CODE_FONT)
    # caption
    add_text(slide, Emu(800000), Emu(5000000), Emu(10600000), Emu(500000),
             "值越接近 1  →  语义越相似", size=20, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_notes(slide, page["notes"])


def render_diagram_4layer(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    add_page_header(slide, "四层记忆架构", idx, total, subtitle="Layer 1 → 4 · 检索优先级")
    # 4 horizontal bands, each with: tag, name, content, store, retention
    bands = [
        ("L1", "工作记忆", "当前对话上下文", "内存", "会话结束清除"),
        ("L2", "短期记忆", "今日对话记录", "MD 文件", "数天到数周"),
        ("L3", "长期记忆", "重要事实 / 偏好", "MEMORY.md + 向量DB", "永久"),
        ("L4", "知识库",   "领域知识 / 技能", "结构化文档 + 索引", "永久（定期更新）"),
    ]
    colors = [
        RGBColor(0xE3, 0xF2, 0xFD),
        RGBColor(0xBB, 0xDE, 0xFB),
        RGBColor(0x90, 0xCA, 0xF9),
        RGBColor(0x64, 0xB5, 0xF6),
    ]
    y0 = Emu(1600000); h = Emu(900000); gap = Emu(80000)
    for i, ((tag, name, content, store, ret), col) in enumerate(zip(bands, colors)):
        yy = y0 + i * (h + gap)
        add_rect(slide, Emu(700000), yy, Emu(10800000), h, col)
        # tag block
        add_rect(slide, Emu(700000), yy, Emu(900000), h, NAVY)
        add_text(slide, Emu(700000), yy, Emu(900000), h, tag,
                 size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # name
        add_text(slide, Emu(1700000), yy, Emu(2200000), h, name,
                 size=20, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        # content
        add_text(slide, Emu(4000000), yy, Emu(3200000), h, content,
                 size=16, color=GRAY_DARK, anchor=MSO_ANCHOR.MIDDLE)
        # store
        add_text(slide, Emu(7300000), yy, Emu(2700000), h, store,
                 size=14, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, font=CODE_FONT)
        # retention
        add_text(slide, Emu(10100000), yy, Emu(1400000), h, ret,
                 size=13, color=ACCENT_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # priority arrow on the left side
    add_text(slide, Emu(140000), Emu(1500000), Emu(500000), Emu(450000),
             "高", size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_arrow_down(slide, Emu(380000), Emu(2050000), Emu(2700000), color=ACCENT, width_pt=2)
    add_text(slide, Emu(140000), Emu(4900000), Emu(500000), Emu(400000),
             "低", size=12, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Emu(50000), Emu(5400000), Emu(680000), Emu(400000),
             "检索\n优先级", size=10, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=1.1)
    add_notes(slide, page["notes"])


def render_diagram_4layer_flow(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    add_page_header(slide, "四层协作检索流程", idx, total,
                    subtitle="一个真实查询如何穿过四层记忆")
    # user question card
    add_round_rect(slide, Emu(700000), Emu(1500000), Emu(10800000), Emu(700000),
                   NAVY)
    add_text(slide, Emu(900000), Emu(1500000), Emu(10400000), Emu(700000),
             "用户：\"我之前让你监控的 Polymarket 市场现在怎样？\"",
             size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # 4 step cards horizontally
    steps = [
        ("步骤 1 · L1", "工作记忆", "当前对话无相关上下文", False),
        ("步骤 2 · L2", "短期记忆", "今日无相关记录", False),
        ("步骤 3 · L3", "长期记忆", "找到：「启动 Polymarket 套利监控」", True),
        ("步骤 4 · L4", "知识库",   "找到：「价格监控脚本每小时扫描」", True),
    ]
    sx0 = Emu(700000); sy = Emu(2500000); sw = Emu(2550000); sh = Emu(2200000); gap = Emu(120000)
    for i, (lbl, name, content, hit) in enumerate(steps):
        x = sx0 + i * (sw + gap)
        bg = RGBColor(0xE6, 0xF7, 0xE6) if hit else RGBColor(0xF2, 0xF2, 0xF2)
        bd = GREEN if hit else GRAY
        card = add_round_rect(slide, x, sy, sw, sh, bg, line_color=bd)
        card.line.width = Pt(1.5)
        add_text(slide, x, sy + Emu(150000), sw, Emu(350000), lbl,
                 size=12, bold=True, color=bd, align=PP_ALIGN.CENTER)
        add_text(slide, x, sy + Emu(550000), sw, Emu(500000), name,
                 size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # divider
        add_rect(slide, x + Emu(400000), sy + Emu(1100000), sw - Emu(800000), Emu(20000),
                 GRAY_LIGHT)
        add_text(slide, x + Emu(150000), sy + Emu(1200000), sw - Emu(300000), Emu(900000),
                 content, size=13, color=GRAY_DARK, align=PP_ALIGN.CENTER, line_spacing=1.4,
                 anchor=MSO_ANCHOR.TOP)
        if hit:
            add_text(slide, x, sy + sh - Emu(380000), sw, Emu(330000),
                     "✓ HIT", size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        else:
            add_text(slide, x, sy + sh - Emu(380000), sw, Emu(330000),
                     "✗ MISS", size=14, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            add_arrow_right(slide, x + sw + Emu(10000), sy + sh // 2,
                            gap - Emu(20000), color=BLUE, width_pt=2)
    # final synthesis
    add_round_rect(slide, Emu(700000), Emu(4900000), Emu(10800000), Emu(700000),
                   ACCENT)
    add_text(slide, Emu(700000), Emu(4900000), Emu(10800000), Emu(700000),
             "→ 综合回复（融合 L3 + L4 多层信息）",
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, page["notes"])


def render_diagram_rrf(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    add_page_header(slide, "RRF 融合算法", idx, total,
                    subtitle="Reciprocal Rank Fusion · 双路稳定者胜出")
    # formula bar
    add_round_rect(slide, Emu(700000), Emu(1550000), Emu(10800000), Emu(750000), NAVY)
    add_text(slide, Emu(700000), Emu(1550000), Emu(10800000), Emu(750000),
             "RRF_score(d) = Σ  1 / (k + rank + 1)        k = 60",
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font=CODE_FONT)
    # 3 doc cards
    docs = [
        ("文档 A", [("BM25", "1"), ("向量", "3")], "1/61 + 1/63", "0.03226", True),
        ("文档 C", [("BM25", "2"), ("向量", "2")], "1/62 + 1/62", "0.03226", True),
        ("文档 D", [("BM25", "1"), ("向量", "—")], "仅 1/61",     "0.01639", False),
    ]
    cx0 = Emu(700000); cy = Emu(2550000); cw = Emu(3550000); ch = Emu(2700000); gap = Emu(180000)
    for i, (name, ranks, expr, score, good) in enumerate(docs):
        x = cx0 + i * (cw + gap)
        bg = LIGHT_BLUE if good else RGBColor(0xFD, 0xEA, 0xEA)
        bd = BLUE if good else RGBColor(0xC0, 0x39, 0x2B)
        card = add_round_rect(slide, x, cy, cw, ch, bg, line_color=bd)
        card.line.width = Pt(1.5)
        add_text(slide, x, cy + Emu(150000), cw, Emu(450000), name,
                 size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # rank rows
        ry = cy + Emu(700000)
        for j, (label, rank) in enumerate(ranks):
            yy = ry + j * Emu(420000)
            add_text(slide, x + Emu(300000), yy, Emu(1400000), Emu(380000),
                     label + "排名", size=14, color=GRAY_DARK)
            add_text(slide, x + cw - Emu(1500000), yy, Emu(1200000), Emu(380000),
                     "#" + rank, size=18, bold=True, color=ACCENT_DARK,
                     align=PP_ALIGN.RIGHT, font=CODE_FONT)
        # expr
        add_text(slide, x, cy + Emu(1700000), cw, Emu(380000),
                 expr, size=14, color=GRAY, align=PP_ALIGN.CENTER, font=CODE_FONT)
        # score
        add_text(slide, x, cy + Emu(2150000), cw, Emu(450000),
                 score, size=28, bold=True,
                 color=GREEN if good else RGBColor(0xC0, 0x39, 0x2B),
                 align=PP_ALIGN.CENTER, font=CODE_FONT)
    # caption
    add_text(slide, Emu(700000), Emu(5450000), Emu(10800000), Emu(550000),
             "RRF 奖励「在多个排名中都表现稳定」的文档 → 单路出现的文档分数被显著压低",
             size=15, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_notes(slide, page["notes"])


# code page
def extract_code_block(body):
    m = re.search(r"```(?:python)?\n(.*?)```", body, re.DOTALL)
    return m.group(1).rstrip() if m else ""


def render_code(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    body = page["body"]
    # find subtitle ### ...
    sub_m = re.search(r"^###\s+(.+)$", body, re.MULTILINE)
    subtitle = sub_m.group(1).strip() if sub_m else None
    add_page_header(slide, subtitle or page["id"], idx, total, subtitle="Python · 代码片段")
    code = extract_code_block(body)
    # caption / trailing text after code
    after = body
    if "```" in after:
        after = re.split(r"```", after)[-1].strip()
    # code container
    cx = Emu(600000); cy = Emu(1500000); cw = Emu(11000000); ch = Emu(4400000)
    bg = add_rect(slide, cx, cy, cw, ch, CODE_BG)
    # subtle top stripe
    add_rect(slide, cx, cy, cw, Emu(60000), ACCENT)
    # window dots
    for i, col in enumerate([RGBColor(0xFF, 0x5F, 0x57),
                             RGBColor(0xFE, 0xBC, 0x2E),
                             RGBColor(0x28, 0xC8, 0x40)]):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     cx + Emu(180000) + i * Emu(280000),
                                     cy + Emu(160000), Emu(180000), Emu(180000))
        dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
    # text frame
    tb = slide.shapes.add_textbox(cx + Emu(120000), cy + Emu(450000),
                                  cw - Emu(240000), ch - Emu(550000))
    render_code_lines(tb.text_frame, code, size=13)
    # caption below
    if after:
        cap = after.replace("**", "")
        add_text(slide, Emu(700000), Emu(6000000), Emu(10800000), Emu(400000),
                 cap, size=14, color=ACCENT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, page["notes"])


def render_table(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    body = page["body"]
    sub_m = re.search(r"^###\s+(.+)$", body, re.MULTILINE)
    title = sub_m.group(1).strip() if sub_m else page["id"]
    add_page_header(slide, title, idx, total, subtitle="对比 · 关键洞察")
    lines = [ln for ln in body.split("\n") if ln.strip().startswith("|")]
    rows = []
    for ln in lines:
        cells = [c.strip().replace("**", "") for c in ln.strip().strip("|").split("|")]
        rows.append(cells)
    data = [r for r in rows[2:] if r and not all(set(c) <= set("-") for c in r)]
    headers = rows[0]
    cols_n = len(headers)
    rows_n = len(data) + 1
    tbl_x = Emu(700000); tbl_y = Emu(1600000)
    tbl_w = Emu(10800000); tbl_h = Emu(4200000)
    table_shape = slide.shapes.add_table(rows_n, cols_n, tbl_x, tbl_y, tbl_w, tbl_h)
    table = table_shape.table
    # column widths: first column narrower
    if cols_n == 2:
        table.columns[0].width = Emu(3800000); table.columns[1].width = Emu(7000000)
    elif cols_n == 3:
        table.columns[0].width = Emu(2800000)
        table.columns[1].width = Emu(4000000)
        table.columns[2].width = Emu(4000000)
    else:
        for c in range(cols_n):
            table.columns[c].width = Emu(10800000 // cols_n)
    # header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.name = CN_FONT; run.font.size = Pt(18); run.font.bold = True
        run.font.color.rgb = WHITE
    for i, row in enumerate(data, start=1):
        for j, txt in enumerate(row[:cols_n]):
            cell = table.cell(i, j)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE if i % 2 == 1 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j != 0 else PP_ALIGN.LEFT
            run = p.add_run(); run.text = txt
            run.font.name = CN_FONT; run.font.size = Pt(15)
            run.font.bold = (j == 0)
            run.font.color.rgb = NAVY if j == 0 else GRAY_DARK
    # caption (lines after table starting with >)
    quote_lines = [ln.strip().lstrip(">").strip() for ln in body.split("\n") if ln.strip().startswith(">")]
    if quote_lines:
        add_text(slide, Emu(700000), Emu(6000000), Emu(10800000), Emu(400000),
                 "  ".join(q.replace("**", "") for q in quote_lines),
                 size=14, color=ACCENT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide, page["notes"])


def render_transition(slide, page, idx, total):
    set_slide_bg(slide, NAVY)
    # accent stripe
    add_rect(slide, Emu(0), Emu(0), Emu(180000), SLIDE_H, ACCENT)
    add_rect(slide, SLIDE_W - Emu(180000), Emu(0), Emu(180000), SLIDE_H, ACCENT)
    # tag
    add_text(slide, Emu(800000), Emu(700000), Emu(4000000), Emu(400000),
             "— 章节过渡 —", size=14, color=ACCENT, bold=True)
    body = page["body"]
    title_m = re.search(r"^#\s+(.+)$", body, re.MULTILINE)
    title = title_m.group(1).strip() if title_m else page["id"]
    add_text(slide, Emu(800000), Emu(1500000), Emu(10500000), Emu(1100000),
             title, size=44, bold=True, color=WHITE)
    rest = body
    if title_m: rest = body[:title_m.start()] + body[title_m.end():]
    rest_lines = [ln.strip() for ln in rest.split("\n") if ln.strip()]
    y = Emu(2900000)
    for ln in rest_lines:
        is_arrow_line = ln.startswith("-->") or "-->" in ln
        clean = ln.replace("**", "")
        if is_arrow_line:
            # split at -->
            parts = clean.split("-->")
            left = parts[0].strip()
            right = parts[-1].strip()
            if left:
                add_text(slide, Emu(800000), y, Emu(10500000), Emu(550000),
                         left, size=22, color=GRAY_LIGHT)
                y += Emu(550000)
            # arrow
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           Emu(800000), y, Emu(900000), Emu(450000))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()
            add_text(slide, Emu(1800000), y - Emu(20000), Emu(9500000), Emu(550000),
                     right, size=24, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
            y += Emu(800000)
        else:
            sz = 24 if "?" in clean or "？" in clean else 22
            color = WHITE if "?" in clean or "？" in clean else GRAY_LIGHT
            bold = "?" in clean or "？" in clean
            add_text(slide, Emu(800000), y, Emu(10500000), Emu(700000),
                     clean, size=sz, bold=bold, color=color, line_spacing=1.4)
            y += Emu(800000)
    # page number
    add_text(slide, Emu(10000000), Emu(6450000), Emu(1700000), Emu(300000),
             f"{idx} / {total}", size=10, color=GRAY_LIGHT, align=PP_ALIGN.RIGHT)
    add_notes(slide, page["notes"])


def render_exercise(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    body = page["body"]
    sub_m = re.search(r"^###\s+(.+)$", body, re.MULTILINE)
    title = sub_m.group(1).strip() if sub_m else page["id"]
    add_page_header(slide, "动手练习", idx, total, subtitle=title)
    # icon + title block
    add_rect(slide, Emu(700000), Emu(1600000), Emu(10800000), Emu(800000), ACCENT)
    add_text(slide, Emu(900000), Emu(1600000), Emu(10500000), Emu(800000),
             "🛠  " + title, size=24, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # parse bullet items "- **key:** value"
    items = []
    for ln in body.split("\n"):
        s = ln.strip()
        if s.startswith("-"):
            content = s.lstrip("-").strip()
            m = re.match(r"\*\*([^*]+)\*\*\s*[:：]\s*(.+)$", content)
            if m:
                items.append((m.group(1), m.group(2)))
            else:
                items.append(("", content.replace("**", "")))
    y0 = Emu(2700000); rh = Emu(650000); gap = Emu(80000)
    for i, (k, v) in enumerate(items):
        yy = y0 + i * (rh + gap)
        bg = LIGHT_BLUE if i % 2 == 0 else WHITE
        add_round_rect(slide, Emu(700000), yy, Emu(10800000), rh, bg, line_color=GRAY_LIGHT)
        add_rect(slide, Emu(700000), yy, Emu(140000), rh, ACCENT)
        if k:
            add_text(slide, Emu(950000), yy, Emu(2000000), rh,
                     k, size=16, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, Emu(3000000), yy, Emu(8400000), rh,
                     v, size=15, color=GRAY_DARK, anchor=MSO_ANCHOR.MIDDLE,
                     font=CODE_FONT if "exercise_" in v or ".py" in v else CN_FONT)
        else:
            add_text(slide, Emu(950000), yy, Emu(10500000), rh,
                     v, size=15, color=GRAY_DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, page["notes"])


def render_summary(slide, page, idx, total):
    set_slide_bg(slide, WHITE)
    add_page_header(slide, "第 1 课  核心总结", idx, total,
                    subtitle="5 条关键要点带回家")
    body = page["body"]
    items = []
    for ln in body.split("\n"):
        m = re.match(r"^\d+\.\s*(.+)$", ln.strip())
        if m:
            items.append(m.group(1).strip())
    y0 = Emu(1700000); rh = Emu(750000); gap = Emu(130000)
    for i, txt in enumerate(items):
        yy = y0 + i * (rh + gap)
        # number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Emu(700000), yy + Emu(120000),
                                        Emu(550000), Emu(550000))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        add_text(slide, Emu(700000), yy + Emu(120000), Emu(550000), Emu(550000),
                 str(i + 1), size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # text card
        add_round_rect(slide, Emu(1400000), yy, Emu(10100000), rh, LIGHT_BLUE)
        # split bold parts
        parts = re.split(r"(\*\*[^*]+\*\*)", txt)
        runs = []
        for p in parts:
            if p.startswith("**") and p.endswith("**"):
                runs.append((p[2:-2], {"size": 18, "bold": True, "color": ACCENT_DARK}))
            else:
                runs.append((p, {"size": 18, "color": NAVY}))
        add_rich_paragraphs(slide, Emu(1700000), yy, Emu(9700000), rh,
                            [{"runs": runs, "align": PP_ALIGN.LEFT}],
                            anchor=MSO_ANCHOR.MIDDLE)
    add_notes(slide, page["notes"])


def render_next_preview(slide, page, idx, total):
    """Final transition page = 'next lesson' preview, special closing styling."""
    set_slide_bg(slide, NAVY)
    if os.path.exists(CLOSING_IMG):
        try:
            slide.shapes.add_picture(CLOSING_IMG, Emu(6800000), Emu(0),
                                     width=Emu(5392000), height=SLIDE_H)
            # left side overlay
            ov = add_rect(slide, Emu(0), Emu(0), Emu(7200000), SLIDE_H, NAVY)
            srgb = ov.fill.fore_color._xFill.find(qn("a:srgbClr"))
            if srgb is not None:
                alpha = etree.SubElement(srgb, qn("a:alpha"))
                alpha.set("val", "92000")
        except Exception as e:
            print("closing image embed failed:", e)
    add_rect(slide, Emu(600000), Emu(900000), Emu(120000), Emu(700000), ACCENT)
    add_text(slide, Emu(800000), Emu(900000), Emu(5500000), Emu(450000),
             "NEXT", size=14, bold=True, color=ACCENT)
    add_text(slide, Emu(800000), Emu(1300000), Emu(6000000), Emu(700000),
             "下节课预告", size=36, bold=True, color=WHITE)
    # body text
    body = page["body"]
    # remove "下节课预告" header line
    lines = []
    for ln in body.split("\n"):
        s = ln.strip()
        if not s: continue
        if s.startswith("#"): continue
        lines.append(s.replace("**", ""))
    y = Emu(2400000)
    for ln in lines:
        if "第2课" in ln:
            add_round_rect(slide, Emu(800000), y, Emu(5800000), Emu(900000), ACCENT)
            add_text(slide, Emu(800000), y, Emu(5800000), Emu(900000),
                     ln, size=24, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            y += Emu(1100000)
        else:
            sz = 22 if ("?" in ln or "？" in ln) else 18
            color = WHITE if ("?" in ln or "？" in ln) else GRAY_LIGHT
            add_text(slide, Emu(800000), y, Emu(5800000), Emu(600000),
                     ln, size=sz, color=color, line_spacing=1.4,
                     bold=("?" in ln or "？" in ln))
            y += Emu(700000)
    # bottom thanks
    add_text(slide, Emu(800000), Emu(6200000), Emu(6000000), Emu(400000),
             "Thanks · 第 1 课结束", size=12, color=GRAY_LIGHT)
    add_text(slide, Emu(10000000), Emu(6450000), Emu(1700000), Emu(300000),
             f"{idx} / {total}", size=10, color=GRAY_LIGHT, align=PP_ALIGN.RIGHT)
    add_notes(slide, page["notes"])


# ===================== DISPATCH =====================

def build():
    pages = parse_slides(SLIDES_MD)
    total = len(pages)
    print(f"parsed {total} pages from slides.md")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for i, page in enumerate(pages, start=1):
        slide = prs.slides.add_slide(blank)
        t = page["type"]
        sid = page["id"]
        try:
            if t == "COVER":
                render_cover(slide, page, i, total)
            elif t == "TOC":
                render_toc(slide, page, i, total)
            elif t == "QUESTION":
                render_question(slide, page, i, total)
            elif t == "CONCEPT":
                render_concept(slide, page, i, total)
            elif t == "DIAGRAM":
                # custom per slide id
                if "1c" in sid:        render_diagram_clawbot_v1(slide, page, i, total)
                elif "2b" in sid:      render_diagram_embedding(slide, page, i, total)
                elif "2d" in sid:      render_diagram_cosine(slide, page, i, total)
                elif "3b" in sid:      render_diagram_4layer(slide, page, i, total)
                elif "3d" in sid:      render_diagram_4layer_flow(slide, page, i, total)
                elif "5b" in sid:      render_diagram_rrf(slide, page, i, total)
                else:
                    render_concept(slide, page, i, total)
            elif t == "CODE":
                render_code(slide, page, i, total)
            elif t == "TABLE":
                render_table(slide, page, i, total)
            elif t == "TRANSITION":
                if "next" in sid:
                    render_next_preview(slide, page, i, total)
                else:
                    render_transition(slide, page, i, total)
            elif t == "EXERCISE":
                render_exercise(slide, page, i, total)
            elif t == "SUMMARY":
                render_summary(slide, page, i, total)
            else:
                render_concept(slide, page, i, total)
        except Exception as e:
            print(f"[ERROR] page {i} ({t} {sid}): {e}")
            raise

    prs.save(OUT)
    # quick stats
    total_shapes = 0
    pics = 0
    for s in prs.slides:
        for sh in s.shapes:
            total_shapes += 1
            if sh.shape_type == 13:
                pics += 1
    print(f"saved -> {OUT}")
    print(f"slides={len(prs.slides)}  total_shapes={total_shapes}  pictures={pics}")


if __name__ == "__main__":
    build()
